"""Standalone indexer
Usage: python index.py
This script builds the Chroma index from files in ./data and persists to ./chroma_db
It is intentionally standalone (does not import Streamlit UI) so you can run indexing in a headless environment.
"""
import os
import io
import sys
from pathlib import Path
from tqdm import tqdm

# Document processing libs
import fitz
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import camelot

# Langchain / Chroma
from langchain.embeddings import OpenAIEmbeddings
from langchain.docstore.document import Document as LCDocument
from langchain.vectorstores import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Read OPENAI_API_KEY from env
if "OPENAI_API_KEY" not in os.environ:
    print("ERROR: Set OPENAI_API_KEY in environment before running.")
    sys.exit(1)


def load_pdf(path):
    doc = fitz.open(path)
    tables_text = ""
    try:
        tables = camelot.read_pdf(path, pages="all")
        for table in tables:
            tables_text += table.df.to_csv(index=False) + "\n"
    except Exception:
        tables_text = ""
    pages = []
    for i, page in enumerate(doc):
        page_text = page.get_text() or ""
        ocr_text = ""
        for img in page.get_images(full=True):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image.get("image")
                if image_bytes:
                    image = Image.open(io.BytesIO(image_bytes))
                    try:
                        ocr_text += pytesseract.image_to_string(image) + "\n"
                    except Exception:
                        pass
            except Exception:
                continue
        content = (page_text + "\n" + ocr_text).strip()
        if i == 0 and tables_text:
            content = content + "\n" + tables_text
        pages.append((os.path.basename(path), content, {"page": i + 1, "type": "pdf"}))
    return pages


def load_docx(path):
    doc = Document(path)
    text = "\n".join(para.text for para in doc.paragraphs)
    tables_text = ""
    for table in doc.tables:
        for row in table.rows:
            row_text = ",".join(cell.text for cell in row.cells)
            tables_text += row_text + "\n"
    content = (text or "") + "\n" + tables_text
    return [(os.path.basename(path), content, {"page": 1, "type": "docx"})]


def load_pptx(path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        tables_text = ""
        ocr_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    row_text = ",".join(cell.text for cell in row.cells)
                    tables_text += row_text + "\n"
            try:
                if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                    image = shape.image
                    image_bytes = image.blob
                    img = Image.open(io.BytesIO(image_bytes))
                    try:
                        ocr_text += pytesseract.image_to_string(img) + "\n"
                    except Exception:
                        pass
            except Exception:
                pass
        content = "\n".join(texts) + "\n" + tables_text + "\n" + ocr_text
        slides.append((os.path.basename(path), content, {"page": i + 1, "type": "pptx"}))
    return slides


def load_all_documents(folder="data"):
    documents = []
    if not os.path.exists(folder):
        print(f"Data folder '{folder}' not found. Create it and add sample files.")
        return documents
    for file in os.listdir(folder):
        full_path = os.path.join(folder, file)
        try:
            if file.lower().endswith(".pdf"):
                documents.extend(load_pdf(full_path))
            elif file.lower().endswith(".docx"):
                documents.extend(load_docx(full_path))
            elif file.lower().endswith(".pptx"):
                documents.extend(load_pptx(full_path))
        except Exception as e:
            print(f"Skipping {file}: {e}")
    return documents


def chunk_documents(docs, chunk_size=700, chunk_overlap=100):
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    all_chunks = []
    for filename, content, meta in docs:
        chunks = splitter.split_text(content or "")
        for idx, chunk in enumerate(chunks):
            metadata = {"source": filename}
            if isinstance(meta, dict):
                metadata.update(meta)
            metadata["chunk_index"] = idx
            all_chunks.append(LCDocument(page_content=chunk, metadata=metadata))
    return all_chunks


def build_chroma_db(chunks, persist_dir="chroma_db"):
    embedding = OpenAIEmbeddings()
    try:
        vectordb = Chroma.from_documents(documents=chunks, embedding=embedding, persist_directory=persist_dir)
        vectordb.persist()
        return vectordb
    except Exception as e:
        print(f"Error building Chroma DB: {e}")
        # delete potentially corrupt directory and try again
        try:
            import shutil, os
            if os.path.isdir(persist_dir):
                shutil.rmtree(persist_dir)
        except Exception:
            pass
        vectordb = Chroma.from_documents(documents=chunks, embedding=embedding, persist_directory=persist_dir)
        vectordb.persist()
        return vectordb


if __name__ == "__main__":
    docs = load_all_documents("data")
    print(f"Loaded {len(docs)} page/slide entries from data/")
    chunks = chunk_documents(docs)
    print(f"Created {len(chunks)} chunks")
    print("Building Chroma DB (this will call OpenAI embeddings)...")
    vectordb = build_chroma_db(chunks)
    print("Index built and persisted to chroma_db/")
