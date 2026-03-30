import streamlit as st
import os
import sys
import warnings
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image
import io
import camelot  # For PDF table extraction
import json
import hashlib
import gc
import datetime
from pathlib import Path
from dotenv import load_dotenv

# Suppress PyMuPDF warnings
warnings.filterwarnings("ignore")
sys.stderr = open(os.devnull, "w")

# Set Tesseract path for Windows (update if needed)
pytesseract.pytesseract.tesseract_cmd = r'C:\Users\Ara86344\AppData\Local\Programs\Tesseract-OCR\tesseract.exe'

from langchain.chat_models import ChatOpenAI
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import Chroma
from langchain.docstore.document import Document as LCDocument
from langchain.chains import RetrievalQA
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Auto-load .env (if present) then read OpenAI API key from env or Streamlit secrets
load_dotenv(dotenv_path=Path(__file__).parent / ".env")
openai_api_key = os.getenv("OPENAI_API_KEY") or (st.secrets["OPENAI_API_KEY"] if hasattr(st, "secrets") and "OPENAI_API_KEY" in st.secrets else None)
if openai_api_key:
        os.environ["OPENAI_API_KEY"] = openai_api_key
else:
    st.warning("OpenAI API key not set. Set OPENAI_API_KEY env var or add to Streamlit secrets.")
    # App can still start but LLM calls will fail until the key is provided




# --- Document Loaders ---
def load_pdf(path):
    doc = fitz.open(path)
    # extract tables once
    tables_text = ""
    try:
        tables = camelot.read_pdf(path, pages="all")
        for table in tables:
            tables_text += table.df.to_string() + "\n"
    except Exception:
        tables_text = ""

    pages = []
    for i, page in enumerate(doc):
        page_text = page.get_text()
        ocr_text = ""
        for img in page.get_images(full=True):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image.get("image")
                if image_bytes:
                    image = Image.open(io.BytesIO(image_bytes))
                    try:
                        ocr_text += pytesseract.image_to_string(image) + "\n"
                    except Exception:
                        pass
            except Exception:
                continue
        content = (page_text or "") + "\n" + ocr_text
        if i == 0 and tables_text:
            content = content + "\n" + tables_text
        pages.append((os.path.basename(path), content, {"page": i + 1, "type": "pdf"}))
    return pages

def load_docx(path):
    doc = Document(path)
    text = "\n".join(para.text for para in doc.paragraphs)
    tables_text = ""
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            tables_text += row_text + "\n"
    content = (text or "") + "\n" + tables_text
    return [(os.path.basename(path), content, {"page": 1, "type": "docx"})]

def load_pptx(path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        tables_text = ""
        ocr_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    tables_text += row_text + "\n"
            try:
                if shape.shape_type == 13 and hasattr(shape, "image"):
                    image = shape.image
                    image_bytes = image.blob
                    img = Image.open(io.BytesIO(image_bytes))
                    try:
                        ocr_text += pytesseract.image_to_string(img) + "\n"
                    except Exception:
                        pass
            except Exception:
                pass
        content = "\n".join(texts) + "\n" + tables_text + "\n" + ocr_text
        slides.append((os.path.basename(path), content, {"page": i + 1, "type": "pptx"}))
    return slides

def load_all_documents(folder):
    documents = []
    for file in os.listdir(folder):
        full_path = os.path.join(folder, file)
        try:
            if file.lower().endswith(".pdf"):
                pages = load_pdf(full_path)
                documents.extend(pages)
            elif file.lower().endswith(".docx"):
                docs = load_docx(full_path)
                documents.extend(docs)
            elif file.lower().endswith(".pptx"):
                slides = load_pptx(full_path)
                documents.extend(slides)
            else:
                continue
        except Exception:
            continue
    return documents

def chunk_documents(docs, chunk_size=700, chunk_overlap=100):
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    all_chunks = []
    for filename, content, meta in docs:
        chunks = splitter.split_text(content or "")
        for idx, chunk in enumerate(chunks):
            metadata = {"source": filename}
            if isinstance(meta, dict):
                metadata.update(meta)
            metadata["chunk_index"] = idx
            all_chunks.append(LCDocument(page_content=chunk, metadata=metadata))
    return all_chunks

def build_chroma_db(chunks, persist_dir="chroma_db"):
    embedding = OpenAIEmbeddings()
    # attempt to create or rebuild the database; on failure try to remove corrupt files and retry
    try:
        vectordb = Chroma.from_documents(documents=chunks, embedding=embedding, persist_directory=persist_dir)
        vectordb.persist()
        return vectordb
    except Exception as e:
        st.error(f"Error building Chroma DB: {e}")
        # remove existing directory and retry once
        try:
            import shutil, os
            if os.path.isdir(persist_dir):
                shutil.rmtree(persist_dir)
        except Exception:
            pass
        # second attempt
        vectordb = Chroma.from_documents(documents=chunks, embedding=embedding, persist_directory=persist_dir)
        vectordb.persist()
        return vectordb

def load_chroma_db(persist_dir="chroma_db"):
    embedding = OpenAIEmbeddings()
    try:
        vectordb = Chroma(persist_directory=persist_dir, embedding_function=embedding)
        return vectordb
    except Exception as e:
        # catch database errors (corruption, compaction failures)
        st.error(f"Failed to load Chroma DB: {e}. The database may be corrupt and will be deleted.")
        import shutil, os
        try:
            if os.path.isdir(persist_dir):
                shutil.rmtree(persist_dir)
        except Exception:
            pass
        return None

# --- Streamlit UI ---
st.set_page_config(page_title="Q Secure AI", page_icon="", layout="wide")
st.title(" Q Secure AI")
# Small CSS to improve chat appearance and ensure good contrast in both light and dark themes
st.markdown(
    """
    <style>
    /* Base styling for chat bubbles */
    .chat-user {background:#e6f2ff;color:#000;padding:12px;border-radius:12px;margin:8px 0;max-width:85%;box-shadow:0 1px 2px rgba(0,0,0,0.06)}
    .chat-bot {background:#ffffff;color:#000;padding:12px;border-radius:12px;margin:8px 0;max-width:85%;box-shadow:0 1px 2px rgba(0,0,0,0.06)}
    .chat-meta {font-size:0.78em;color:#666;margin-bottom:6px}
    /* Slightly different alignment for user vs bot */
    .chat-user {margin-left:auto}
    .chat-bot {margin-right:auto}
    /* Make sure text wraps nicely */
    .chat-user, .chat-bot {white-space:pre-wrap;word-break:break-word}

    /* Dark theme adjustments */
    @media (prefers-color-scheme: dark) {
        .chat-user {background:#2b4156;color:#e6f4ff;box-shadow:0 1px 2px rgba(0,0,0,0.6)}
        .chat-bot {background:#16202b;color:#e6eef6;box-shadow:0 1px 2px rgba(0,0,0,0.6)}
        .chat-meta {color:#9aa6b2}
    }

    /* Small responsive tweak */
    @media (max-width: 600px) {
        .chat-user, .chat-bot {max-width:100%}
    }
    </style>
    """,
    unsafe_allow_html=True,
)

# One-time document processing
def process_documents():
    # Process documents file-by-file so we can show progress in the UI.
    data_dir = "data"
    files = [f for f in os.listdir(data_dir) if f.lower().endswith((".pdf", ".docx", ".pptx"))]
    total_files = max(1, len(files))
    progress = st.progress(0)
    processed_docs = []
    idx = 0
    with st.spinner("Processing documents and building index. This may take a while..."):
        for fname in files:
            idx += 1
            full_path = os.path.join(data_dir, fname)
            try:
                if fname.lower().endswith(".pdf"):
                    entries = load_pdf(full_path)
                elif fname.lower().endswith(".docx"):
                    entries = load_docx(full_path)
                elif fname.lower().endswith(".pptx"):
                    entries = load_pptx(full_path)
                else:
                    entries = []
                if entries:
                    processed_docs.extend(entries)
            except Exception:
                # skip problematic file but keep processing
                pass
            progress.progress(int(idx / total_files * 100))
        # chunk all collected entries and build the vectordb once
        chunks = chunk_documents(processed_docs)
        vectordb = build_chroma_db(chunks)
    progress.empty()
    return vectordb


MANIFEST_PATH = "chroma_db/manifest.json"


def compute_file_hash(path):
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def load_manifest():
    try:
        with open(MANIFEST_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def save_manifest(manifest):
    os.makedirs(os.path.dirname(MANIFEST_PATH), exist_ok=True)
    with open(MANIFEST_PATH, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2)


def incremental_index(vectordb, data_folder="data"):
    """Index only new or changed files in data_folder into existing vectordb.
    Returns number of documents (chunks) added.
    """
    manifest = load_manifest()
    added = 0
    files = [f for f in os.listdir(data_folder) if f.lower().endswith((".pdf", ".docx", ".pptx"))]
    for fname in files:
        full_path = os.path.join(data_folder, fname)
        try:
            file_hash = compute_file_hash(full_path)
        except Exception:
            continue
        prev_hash = manifest.get(fname, {}).get("hash")
        if prev_hash == file_hash:
            # already indexed
            continue
        # load and chunk this file only
        try:
            entries = []
            if fname.lower().endswith(".pdf"):
                entries = load_pdf(full_path)
            elif fname.lower().endswith(".docx"):
                entries = load_docx(full_path)
            elif fname.lower().endswith(".pptx"):
                entries = load_pptx(full_path)
            chunks = chunk_documents(entries)
            if chunks:
                # add to existing vectordb
                try:
                    vectordb.add_documents(documents=chunks)
                except Exception:
                    # fallback: rebuild partial DB via from_documents
                    vectordb = Chroma.from_documents(documents=chunks, embedding=OpenAIEmbeddings(), persist_directory="chroma_db")
                vectordb.persist()
                added += len(chunks)
            # update manifest
            manifest[fname] = {"hash": file_hash, "last_indexed": int(os.path.getmtime(full_path))}
            save_manifest(manifest)
        except Exception:
            continue
    return added


def show_chroma_diagnostics(path="chroma_db"):
    """Top-level diagnostics for the chroma DB usable from the sidebar or elsewhere."""
    if not os.path.exists(path):
        st.info("No chroma_db directory present.")
        return
    total_size = 0
    total_files = 0
    for root, dirs, files in os.walk(path):
        for f in files:
            fp = os.path.join(root, f)
            try:
                total_size += os.path.getsize(fp)
                total_files += 1
            except Exception:
                pass
    st.write(f"Files: {total_files}")
    st.write(f"Total size: {total_size / 1024:.1f} KB")
    # Try to show manifest if present
    manifest_path = os.path.join(path, "manifest.json")
    if os.path.exists(manifest_path):
        try:
            with open(manifest_path, "r", encoding="utf-8") as f:
                manifest = json.load(f)
            st.write("Manifest (sample):")
            sample_items = dict(list(manifest.items())[:10])
            st.json(sample_items)
        except Exception:
            st.write("Could not read manifest.json")
    else:
        st.write("No manifest.json found in chroma_db")

import shutil
import gc
import time
# Load or build ChromaDB only once
if "vectordb" not in st.session_state:
    if not os.path.exists("chroma_db") or not os.listdir("chroma_db"):
        st.info("Processing and embedding documents. Please wait...")
        st.session_state.vectordb = process_documents()
        st.success("Documents processed and indexed!")
    else:
        st.session_state.vectordb = load_chroma_db()

# Create QA chain only once
if "qa_chain" not in st.session_state:
    retriever = st.session_state.vectordb.as_retriever(search_kwargs={"k": 4})
    llm = ChatOpenAI(temperature=0, model="gpt-4o")  # Latest ChatGPT model
    st.session_state.qa_chain = RetrievalQA.from_chain_type(
        llm=llm, retriever=retriever, return_source_documents=True
    )

# --- Chat UI ---
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

def chat(query):
    # return LLM response and similarity scores from Chroma
    vect = st.session_state.vectordb
    try:
        docs_scores = vect.similarity_search_with_score(query, k=4)
    except Exception:
        docs_scores = []
    response = st.session_state.qa_chain(query)
    return response, docs_scores

with st.form("chat_form"):
    user_input = st.text_input("Ask a question:", "")
    submitted = st.form_submit_button("Send")
    if submitted and user_input:
        response, docs_scores = chat(user_input)
        st.session_state.chat_history.append(
            {"user": user_input, "bot": response["result"], "sources": response["source_documents"], "scores": docs_scores, "time": datetime.datetime.now().isoformat()}
        )

# Display chat history with nicer layout
for msg in st.session_state.chat_history:
    # User message
    st.markdown(f"<div class='chat-meta'>You · {msg.get('time','')}</div>", unsafe_allow_html=True)
    st.markdown(f"<div class='chat-user'>{msg['user']}</div>", unsafe_allow_html=True)
    # Bot message
    st.markdown(f"<div class='chat-meta'>Gen AI </div>", unsafe_allow_html=True)
    st.markdown(f"<div class='chat-bot'>{msg['bot']}</div>", unsafe_allow_html=True)
    # Sources in an expander (deduplicated)
    with st.expander("Source documents (click to expand)"):
        seen_sources = set()
        for doc in msg.get("sources", []):
            src = getattr(doc, 'metadata', {}).get('source') if hasattr(doc, 'metadata') else doc.metadata.get('source') if doc else None
            if not src or src in seen_sources:
                continue
            seen_sources.add(src)
            page = getattr(doc, 'metadata', {}).get('page') if hasattr(doc, 'metadata') else doc.metadata.get('page') if doc else None
            if page:
                st.write(f"- {src} (page {page})")
            else:
                st.write(f"- {src}")
    st.markdown("---")

# Sidebar controls: re-index and upload files
with st.sidebar:
    st.header("Controls")
    if st.button("Re-index documents"):
        # Helper: try to cleanly remove chroma_db with retries; if removal fails,
        # we'll fall back to incremental indexing using the existing vectordb.
        def safe_remove_dir(path, max_attempts=6, base_delay=0.5):
            """Try to remove a directory with retries and on-error handler to
            fix permission bits on Windows. Returns True if removed, False otherwise."""
            if not os.path.exists(path):
                return True

            # Try to release known references/clients
            try:
                if "vectordb" in st.session_state and st.session_state.vectordb is not None:
                    try:
                        if hasattr(st.session_state.vectordb, "persist"):
                            st.session_state.vectordb.persist()
                    except Exception:
                        pass
                    client = getattr(st.session_state.vectordb, "_client", None) or getattr(st.session_state.vectordb, "client", None)
                    if client is not None:
                        for m in ("close", "shutdown", "disconnect", "stop"):
                            if hasattr(client, m):
                                try:
                                    getattr(client, m)()
                                except Exception:
                                    pass
            except Exception:
                pass

            # Remove references and force GC so OS locks (Windows) are more likely released
            try:
                if "vectordb" in st.session_state:
                    del st.session_state["vectordb"]
            except Exception:
                pass
            try:
                if "qa_chain" in st.session_state:
                    del st.session_state["qa_chain"]
            except Exception:
                pass
            gc.collect()

            def _on_rm_error(func, path_inner, exc_info):
                import stat
                try:
                    os.chmod(path_inner, stat.S_IWRITE)
                    func(path_inner)
                except Exception:
                    pass

            attempt = 0
            while attempt < max_attempts:
                try:
                    # small pause to allow OS to release handles
                    time.sleep(base_delay * (1 + attempt * 0.5))
                    shutil.rmtree(path, onerror=_on_rm_error)
                    return True
                except Exception as e:
                    attempt += 1
                    st.warning(f"Attempt {attempt}/{max_attempts} to remove {path} failed: {e}")
                    gc.collect()
                    time.sleep(base_delay * (1 + attempt * 0.5))
            return False

        def show_chroma_diagnostics(path="chroma_db"):
            st.subheader("Chroma DB diagnostics")
            if not os.path.exists(path):
                st.info("No chroma_db directory present.")
                return
            total_size = 0
            total_files = 0
            for root, dirs, files in os.walk(path):
                for f in files:
                    fp = os.path.join(root, f)
                    try:
                        total_size += os.path.getsize(fp)
                        total_files += 1
                    except Exception:
                        pass
            st.write(f"Files: {total_files}")
            st.write(f"Total size: {total_size / 1024:.1f} KB")
            # Try to show manifest if present
            manifest_path = os.path.join(path, "manifest.json")
            if os.path.exists(manifest_path):
                try:
                    with open(manifest_path, "r", encoding="utf-8") as f:
                        manifest = json.load(f)
                    st.write("Manifest (sample):")
                    sample_items = dict(list(manifest.items())[:10])
                    st.json(sample_items)
                except Exception:
                    st.write("Could not read manifest.json")
            else:
                st.write("No manifest.json found in chroma_db")

        # Attempt safe removal
        removed = safe_remove_dir("chroma_db")
        if not removed:
            st.warning("Could not fully remove chroma_db directory. Falling back to incremental indexing where possible.")
            show_chroma_diagnostics("chroma_db")

        # Run re-index in-place and rebuild vectordb and QA chain
        with st.spinner("Re-indexing documents now. This may take a few minutes..."):
            try:
                # If chroma_db still exists and a vectordb is available in session, try incremental indexing
                if os.path.exists("chroma_db") and "vectordb" in st.session_state and st.session_state.vectordb is not None:
                    added = incremental_index(st.session_state.vectordb)
                    # recreate retriever/qa chain
                    retriever = st.session_state.vectordb.as_retriever(search_kwargs={"k": 4})
                    llm = ChatOpenAI(temperature=0, model="gpt-4o")
                    st.session_state.qa_chain = RetrievalQA.from_chain_type(llm=llm, retriever=retriever, return_source_documents=True)
                    st.success(f"Re-index complete. {added} new chunks indexed (incremental).")
                else:
                    # No existing vectordb (or chroma_db removed): do a full processing
                    new_vectordb = process_documents()
                    st.session_state.vectordb = new_vectordb
                    retriever = st.session_state.vectordb.as_retriever(search_kwargs={"k": 4})
                    llm = ChatOpenAI(temperature=0, model="gpt-4o")
                    st.session_state.qa_chain = RetrievalQA.from_chain_type(llm=llm, retriever=retriever, return_source_documents=True)
                    st.success("Full re-index complete. Ready to answer questions.")
            except Exception as e:
                st.error(f"Re-index failed: {e}")

    uploaded_files = st.file_uploader("Upload documents (PDF/DOCX/PPTX)", type=["pdf", "docx", "pptx"], accept_multiple_files=True)
    # Avoid saving the same uploaded files multiple times in a session.
    last_names = st.session_state.get("uploads_last_names")
    uploaded_names = [up.name for up in uploaded_files] if uploaded_files else []
    if uploaded_files and uploaded_names != last_names:
        os.makedirs("data", exist_ok=True)
        saved = 0
        for up in uploaded_files:
            safe_name = os.path.basename(up.name)
            out_path = os.path.join("data", safe_name)
            file_bytes = up.getbuffer()
            try:
                if os.path.exists(out_path):
                    # compare contents to avoid creating timestamped duplicates
                    try:
                        with open(out_path, "rb") as existing:
                            existing_bytes = existing.read()
                        if existing_bytes == file_bytes:
                            # same file already present, skip
                            st.info(f"File already exists and is identical: {safe_name}")
                            continue
                    except Exception:
                        pass
                    # different content, write with timestamp suffix
                    alt_path = os.path.join("data", f"{int(time.time())}_{safe_name}")
                    with open(alt_path, "wb") as f:
                        f.write(file_bytes)
                    saved += 1
                else:
                    # write new file
                    with open(out_path, "wb") as f:
                        f.write(file_bytes)
                    saved += 1
            except PermissionError:
                # fallback to uploads subfolder
                alt_dir = os.path.join("data", "uploads")
                os.makedirs(alt_dir, exist_ok=True)
                alt_path = os.path.join(alt_dir, f"{int(time.time())}_{safe_name}")
                try:
                    with open(alt_path, "wb") as f:
                        f.write(file_bytes)
                    saved += 1
                    st.warning(f"Saved {safe_name} to alternate path due to permission restrictions: {alt_path}")
                except Exception as e:
                    st.error(f"Failed to save {safe_name}: {e}")
            except Exception as e:
                st.error(f"Failed to save {safe_name}: {e}")
        if saved:
            st.success(f"Saved {saved} files to data/. Click Re-index to process them.")
        # remember which upload names we've saved so we don't re-save them
        st.session_state["uploads_last_names"] = uploaded_names

    # Sidebar diagnostics quick view
    with st.expander("Chroma DB diagnostics"):
        show_chroma_diagnostics("chroma_db")
